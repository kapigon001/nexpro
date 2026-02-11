"""
ネクプロ全社戦略プレゼンテーション pptx生成スクリプト
株主総会・経営会議品質 — 白ベース×紺/グレー、游ゴシック指定
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE
import os

# ==================================================================
# Constants
# ==================================================================
CHART_DIR = '/home/user/nexpro/chart_images'
OUT_PATH = '/home/user/nexpro/nexpro_strategy_presentation.pptx'

FONT_JP = '游ゴシック'
FONT_EN = 'Yu Gothic'

NAVY = RGBColor(0x1B, 0x2A, 0x4A)
BLUE = RGBColor(0x2E, 0x6D, 0xA4)
LIGHT_BLUE = RGBColor(0x5B, 0xA4, 0xCF)
ACCENT = RGBColor(0xE8, 0x91, 0x3A)
ACCENT_RED = RGBColor(0xD6, 0x40, 0x45)
GREY = RGBColor(0x8C, 0x8C, 0x8C)
DARK_GREY = RGBColor(0x4A, 0x4A, 0x4A)
LIGHT_GREY = RGBColor(0xD9, 0xD9, 0xD9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG_WHITE = RGBColor(0xFA, 0xFA, 0xFA)
TABLE_HEADER_BG = RGBColor(0x1B, 0x2A, 0x4A)
TABLE_ROW_ALT = RGBColor(0xF0, 0xF4, 0xF8)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H


# ==================================================================
# Helper Functions
# ==================================================================
def add_slide():
    """Add a blank slide."""
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    # White background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    return slide


def set_font(run, size=12, bold=False, italic=False, color=NAVY, name=None):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = name or FONT_JP


def add_textbox(slide, left, top, width, height, text='', size=12, bold=False,
                color=NAVY, alignment=PP_ALIGN.LEFT, name=None, italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    set_font(run, size=size, bold=bold, color=color, name=name, italic=italic)
    return txBox


def add_multiline_textbox(slide, left, top, width, height, lines,
                          default_size=11, default_color=NAVY, line_spacing=1.2):
    """lines: list of (text, size, bold, color) or just str."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(2)

        if isinstance(line, str):
            run = p.add_run()
            run.text = line
            set_font(run, size=default_size, color=default_color)
        elif isinstance(line, tuple):
            text = line[0]
            size = line[1] if len(line) > 1 else default_size
            bold = line[2] if len(line) > 2 else False
            color = line[3] if len(line) > 3 else default_color
            run = p.add_run()
            run.text = text
            set_font(run, size=size, bold=bold, color=color)
    return txBox


def add_accent_bar(slide, left, top, width=Inches(0.08), height=Inches(0.6), color=ACCENT):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_bottom_bar(slide):
    """Add bottom navy bar with copyright."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0), Inches(7.0), SLIDE_W, Inches(0.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    tf = bar.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = 'CONFIDENTIAL  |  NexPro Inc.  |  2026'
    set_font(run, size=8, color=WHITE)


def add_slide_number(slide, num, total=20):
    add_textbox(slide, Inches(12.5), Inches(7.05), Inches(0.7), Inches(0.35),
                f'{num}/{total}', size=8, color=WHITE, alignment=PP_ALIGN.RIGHT)


def add_header(slide, title, subtitle=None, slide_num=None):
    """Standard slide header with accent bar."""
    add_accent_bar(slide, Inches(0.6), Inches(0.4), Inches(0.07), Inches(0.55))
    add_textbox(slide, Inches(0.8), Inches(0.35), Inches(10), Inches(0.6),
                title, size=24, bold=True, color=NAVY)
    if subtitle:
        add_textbox(slide, Inches(0.85), Inches(0.95), Inches(10), Inches(0.4),
                    subtitle, size=13, bold=False, color=GREY, italic=True)
    add_bottom_bar(slide)
    if slide_num:
        add_slide_number(slide, slide_num)


def add_key_message_box(slide, text, left=Inches(0.8), top=Inches(1.4),
                        width=Inches(11.5), height=Inches(0.65)):
    """Blue-bordered key message box."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xE8, 0xF4, 0xFD)
    shape.line.color.rgb = BLUE
    shape.line.width = Pt(1.5)

    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = f'  Key Message:  {text}'
    set_font(run, size=12, bold=True, color=NAVY)


def add_image(slide, img_name, left, top, width=None, height=None):
    path = os.path.join(CHART_DIR, img_name)
    if os.path.exists(path):
        kwargs = {'left': left, 'top': top}
        if width:
            kwargs['width'] = width
        if height:
            kwargs['height'] = height
        slide.shapes.add_picture(path, **kwargs)


def make_table(slide, left, top, width, height, rows, cols, data,
               col_widths=None, header_color=TABLE_HEADER_BG, font_size=9):
    """Create a styled table. data is list of lists."""
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = ''
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
            run = p.add_run()
            run.text = str(data[r][c]) if r < len(data) and c < len(data[r]) else ''

            if r == 0:  # header
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
                set_font(run, size=font_size, bold=True, color=WHITE)
            else:
                if r % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = TABLE_ROW_ALT
                else:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = WHITE
                set_font(run, size=font_size, bold=False, color=DARK_GREY)

            # vertical alignment
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            # margins
            cell.margin_left = Pt(4)
            cell.margin_right = Pt(4)
            cell.margin_top = Pt(2)
            cell.margin_bottom = Pt(2)

    return table_shape


# ==================================================================
# SLIDE 1: Title
# ==================================================================
def slide_01_title():
    slide = add_slide()
    # Navy header band
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0), Inches(0), SLIDE_W, Inches(3.2))
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.line.fill.background()

    add_textbox(slide, Inches(0.8), Inches(0.8), Inches(11), Inches(1.0),
                'ネクプロ 全社戦略提案', size=36, bold=True, color=WHITE)
    add_textbox(slide, Inches(0.8), Inches(1.8), Inches(11), Inches(0.6),
                'ウェビナーツールから B2B エンゲージメント・インテリジェンス基盤へ',
                size=18, color=RGBColor(0xAA, 0xCC, 0xEE))
    add_textbox(slide, Inches(0.8), Inches(2.5), Inches(5), Inches(0.4),
                '対象期間：2026-2028（0-36ヶ月）  |  経営層・マネージャーMTG',
                size=12, color=RGBColor(0x88, 0xAA, 0xCC))

    # 3 decision boxes
    decisions = [
        ('Decision 1', '資源配分の決定', 'プロダクト45% / GTM35% / 組織20%'),
        ('Decision 2', 'ポジショニング転換', 'ウェビナーツール→エンゲージメント基盤'),
        ('Decision 3', '組織再編の承認', 'PMM/RevOps新設・KPIオーナー制度'),
    ]
    for i, (label, title, desc) in enumerate(decisions):
        x = Inches(0.8 + i * 4.1)
        y = Inches(3.8)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      x, y, Inches(3.7), Inches(1.8))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0xF5, 0xF7, 0xFA)
        box.line.color.rgb = LIGHT_BLUE
        box.line.width = Pt(1)

        # accent top bar
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      x, y, Inches(3.7), Inches(0.06))
        bar.fill.solid()
        bar.fill.fore_color.rgb = [ACCENT, BLUE, LIGHT_BLUE][i]
        bar.line.fill.background()

        add_textbox(slide, x + Inches(0.2), y + Inches(0.2), Inches(3.3), Inches(0.3),
                    label, size=10, bold=True, color=ACCENT)
        add_textbox(slide, x + Inches(0.2), y + Inches(0.55), Inches(3.3), Inches(0.4),
                    title, size=16, bold=True, color=NAVY)
        add_textbox(slide, x + Inches(0.2), y + Inches(1.1), Inches(3.3), Inches(0.5),
                    desc, size=10, color=GREY)

    add_bottom_bar(slide)
    add_textbox(slide, Inches(0.8), Inches(6.3), Inches(5), Inches(0.3),
                'CONFIDENTIAL — 社外秘', size=10, bold=True, color=ACCENT_RED)


# ==================================================================
# SLIDE 2: Executive Summary
# ==================================================================
def slide_02_exec_summary():
    slide = add_slide()
    add_header(slide, 'エグゼクティブサマリー',
               '結論 — 「配信ツール」に留まれば死ぬ。「データ基盤」に転換すれば勝てる。', 1)

    add_key_message_box(slide,
        'ネクプロは「ウェビナーツール」→「B2Bエンゲージメント・インテリジェンス基盤」への転換を今すぐ決断すべき')

    # Summary table
    data = [
        ['項目', '現状（FY24実績）', '目標（FY27計画）', '変化率'],
        ['事業定義', 'ウェビナー配信ツール', 'エンゲージメント基盤', '—'],
        ['売上', '¥512M', '¥1,382M', '+170%'],
        ['MRR（年間）', '¥287M', '¥526M', '+83%'],
        ['ARPA（長期PF）', '¥148K/月', '¥204K/月', '+38%'],
        ['月次解約率', '1.7%', '1.0%（目標）', '−0.7pt'],
        ['長期PFアカウント数', '167社', '210社', '+26%'],
        ['新収益柱', 'なし', '¥378M（全体27%）', '—'],
    ]
    make_table(slide, Inches(0.8), Inches(2.2), Inches(11.5), Inches(3.2),
               8, 4, data, font_size=10,
               col_widths=[Inches(2.5), Inches(3.2), Inches(3.5), Inches(2.3)])

    # Bottom insight
    add_multiline_textbox(slide, Inches(0.8), Inches(5.6), Inches(11.5), Inches(1.2), [
        ('根拠: ', 11, True, NAVY),
        ('AIエージェント時代にUI中心SaaSの中間層は圧縮 / 日本ウェビナーSaaSはAI統合で12-18ヶ月遅れ → 先行者優位の猶予あり', 10, False, DARK_GREY),
    ])
    add_multiline_textbox(slide, Inches(0.8), Inches(6.1), Inches(11.5), Inches(0.8), [
        ('反証リスク: ', 11, True, ACCENT_RED),
        ('30名体制での転換実行力 / 新収益柱（営業DX ¥327M）は未実証の仮説値', 10, False, DARK_GREY),
    ])


# ==================================================================
# SLIDE 3: External Environment
# ==================================================================
def slide_03_external():
    slide = add_slide()
    add_header(slide, '外部環境：AIエージェント時代の産業変化',
               'SaaS中間層の圧縮 — AIエージェントが変えるソフトウェア産業の構造', 2)
    add_key_message_box(slide,
        'AIエージェントはSaaSの「UIレイヤー」を不要にするが、「データレイヤー」の価値は増大する')

    add_image(slide, 'saas_layers.png', Inches(0.5), Inches(2.2), width=Inches(6.5))

    # Key facts on right side
    facts = [
        ('$2,850億（42兆円）', 'SaaS株時価総額消失\n（2026.2 Claude Cowork発表後）'),
        ('−43%', 'バーティカルSaaS年初来下落率'),
        ('70%', 'シートベース→成果ベース課金移行\n（2028年まで、IDC予測）'),
        ('4%', '日本のSaaS浸透率（米国15-18%）\n→ 構造的参入障壁＝時間稼ぎ'),
    ]
    for i, (num, desc) in enumerate(facts):
        y = Inches(2.3 + i * 1.15)
        add_textbox(slide, Inches(7.3), y, Inches(2), Inches(0.4),
                    num, size=22, bold=True, color=ACCENT)
        add_textbox(slide, Inches(7.3), y + Inches(0.4), Inches(5.5), Inches(0.6),
                    desc, size=9, color=DARK_GREY)

    add_multiline_textbox(slide, Inches(0.8), Inches(6.3), Inches(11.5), Inches(0.6), [
        ('ネクプロへの示唆: ', 10, True, NAVY),
        ('配信機能 = battlegrounds（蚕食領域） / エンゲージメントデータ = gold mines（価値増大領域）', 10, False, DARK_GREY),
    ])


# ==================================================================
# SLIDE 4: Current Status - Revenue
# ==================================================================
def slide_04_current_revenue():
    slide = add_slide()
    add_header(slide, '自社現状①：売上推移と構成',
               'MRRは安定するも成長率は鈍化。新収益柱の立ち上げが急務。', 3)
    add_key_message_box(slide,
        'FY24成長率+3.1%に急減速。¥1B突破には新収益柱（営業DX・コンパウンド）の成功が不可欠')

    add_image(slide, 'revenue_trend.png', Inches(0.5), Inches(2.2), width=Inches(7))

    # Revenue composition callout
    add_multiline_textbox(slide, Inches(7.8), Inches(2.3), Inches(5), Inches(4.0), [
        ('FY27売上構成（計画）', 13, True, NAVY),
        ('', 6, False, NAVY),
        ('MRR（システム利用料）', 11, True, BLUE),
        ('  ¥527M（38%）', 11, False, DARK_GREY),
        ('', 4, False, NAVY),
        ('オプションサービス', 11, True, LIGHT_BLUE),
        ('  ¥478M（35%）', 11, False, DARK_GREY),
        ('', 4, False, NAVY),
        ('営業DX（新規）', 11, True, ACCENT),
        ('  ¥327M（24%） ← 未実証・最大リスク', 11, False, ACCENT_RED),
        ('', 4, False, NAVY),
        ('コンパウンド（新規）', 11, True, ACCENT),
        ('  ¥51M（4%）', 11, False, DARK_GREY),
        ('', 8, False, NAVY),
        ('※ 営業DXがFY27計画の24%を占める。', 9, True, ACCENT_RED),
        ('  この実行リスクが全計画の成否を左右。', 9, False, ACCENT_RED),
    ])


# ==================================================================
# SLIDE 5: Current Status - KPIs
# ==================================================================
def slide_05_current_kpis():
    slide = add_slide()
    add_header(slide, '自社現状②：SaaS KPI分析',
               'ARPA改善は順調だが、解約率は依然「バケツの穴」状態', 4)

    add_image(slide, 'mrr_arpa.png', Inches(0.3), Inches(1.5), width=Inches(6.2))
    add_image(slide, 'churn_rate.png', Inches(6.5), Inches(1.5), width=Inches(6.2))

    # Key KPI table
    data = [
        ['KPI', 'FY23', 'FY24', 'FY25計画', '評価'],
        ['長期PFアカウント数', '151社', '167社', '179社', '○ 回復傾向'],
        ['ARPA長期PF', '¥138K', '¥148K', '¥169K', '○ 改善中'],
        ['月次解約率（長期）', '2.3%', '1.7%', '1.0%', '△ 要改善'],
        ['新規成約数/年', '27社', '50社', '38社', '○ FY24回復'],
        ['成約率', '9.5%', '10.5%', '—', '△ 業界並み'],
        ['年換算解約率', '24.5%', '18.5%', '11.3%', '× SaaS優良=5%'],
    ]
    make_table(slide, Inches(0.5), Inches(5.0), Inches(12.3), Inches(2.0),
               7, 5, data, font_size=9,
               col_widths=[Inches(3.0), Inches(1.8), Inches(1.8), Inches(1.8), Inches(3.9)])


# ==================================================================
# SLIDE 6: Account & New Revenue
# ==================================================================
def slide_06_accounts_new_rev():
    slide = add_slide()
    add_header(slide, '自社現状③：顧客基盤と新収益柱',
               'アカウント数は回復基調。新規事業は急成長を前提とした計画。', 5)

    add_image(slide, 'accounts.png', Inches(0.3), Inches(1.5), width=Inches(6.2))
    add_image(slide, 'new_revenue.png', Inches(6.5), Inches(1.5), width=Inches(6.5))

    add_multiline_textbox(slide, Inches(0.8), Inches(5.8), Inches(11.5), Inches(1.0), [
        ('課題仮説: ', 11, True, NAVY),
        ('① 利用深度のばらつき（上位/下位の格差大）  ', 10, False, DARK_GREY),
        ('② 受注後オンボーディングの再現性不足  ', 10, False, DARK_GREY),
        ('③ 高単価化のストーリー不足（データ活用→成果証明の接続が弱い）', 10, False, DARK_GREY),
    ])


# ==================================================================
# SLIDE 7: SWOT
# ==================================================================
def slide_07_swot():
    slide = add_slide()
    add_header(slide, '戦略的SWOT分析',
               '「データモート」が唯一の持続的優位性', 6)

    # 4 quadrants
    quad_data = [
        ('Strengths', Inches(0.5), Inches(1.4), NAVY,
         ['ITreview 13期連続受賞の製品品質',
          'Salesforce深度統合（Agentforce国内最先行）',
          '400-500社のエンタープライズ顧客基盤',
          '日本語運用ノウハウ・伴走支援体制',
          '擬似ライブ配信・メディアサイトビルダー']),
        ('Weaknesses', Inches(6.7), Inches(1.4), ACCENT_RED,
         ['30名の少数体制（開発速度限界）',
          '総調達額¥7.5億（資金格差大）',
          'ブランド認知度（Zoom/V-CUBEに劣後）',
          'インテントデータ未整備',
          'PMM/RevOps機能の不在']),
        ('Opportunities', Inches(0.5), Inches(4.2), BLUE,
         ['日本ウェビナーSaaSのAI統合12-18ヶ月遅れ',
          'ON24 Cvent買収→日本優先度低下',
          'Cookie廃止→1stパーティデータ価値増',
          '日本B2B DX余地（SaaS浸透率4%）',
          '企業あたりウェビナー数急増（13→47回/年）']),
        ('Threats', Inches(6.7), Inches(4.2), DARK_GREY,
         ['Zoom/Webex/TeamsのAIエージェント搭載',
          'AIによる配信機能のコモディティ化',
          'Cvent大型統合（ON24+Goldcast ¥700億）',
          '価格競争の激化',
          '汎用AIで差別化希薄化']),
    ]

    for title, x, y, color, items in quad_data:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      x, y, Inches(6.0), Inches(2.6))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0xF8, 0xF9, 0xFB)
        box.line.color.rgb = color
        box.line.width = Pt(1.5)

        # Title bar
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      x, y, Inches(6.0), Inches(0.4))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        add_textbox(slide, x + Inches(0.15), y + Inches(0.02), Inches(5.5), Inches(0.35),
                    title, size=13, bold=True, color=WHITE)

        for j, item in enumerate(items):
            add_textbox(slide, x + Inches(0.2), y + Inches(0.5 + j * 0.38),
                        Inches(5.5), Inches(0.35),
                        f'• {item}', size=9, color=DARK_GREY)


# ==================================================================
# SLIDE 8: Competitive Comparison Table
# ==================================================================
def slide_08_competitive_table():
    slide = add_slide()
    add_header(slide, '競合比較表',
               '5社×7軸の定量・定性比較', 7)
    add_key_message_box(slide,
        'グローバル勢は日本適合が弱く、国内勢はデータ活用が弱い → 両方持てば唯一の存在')

    data = [
        ['評価軸', 'ネクプロ', 'Zoom', 'ON24', 'EventHub', 'bizibl', 'FanGrowth'],
        ['配信品質', '◎', '◎', '○', '○', '○', '△'],
        ['MA/CRM連携', '◎ SF深度', '○ API', '◎ Eloqua', '△', '△', '△'],
        ['日本語適合', '◎', '△', '△', '◎', '◎', '◎'],
        ['データ分析', '○ (発展余地)', '△', '◎ 40-50pt', '△', '△', '△'],
        ['AI機能', '○ Agentforce', '○ AI Comp.', '◎ ACE AI', '△', '△', '△'],
        ['価格柔軟性', '◎', '○', '△ 高額', '○', '◎', '◎'],
        ['支援体制', '◎ 伴走型', '△ セルフ', '○ 海外', '○', '○', '◎ 運用代行'],
    ]
    make_table(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(3.8),
               8, 7, data, font_size=9,
               col_widths=[Inches(1.8), Inches(2.0), Inches(1.6), Inches(2.0),
                           Inches(1.6), Inches(1.6), Inches(1.7)])

    add_multiline_textbox(slide, Inches(0.8), Inches(6.3), Inches(11.5), Inches(0.5), [
        ('示唆: ', 10, True, NAVY),
        ('ネクプロは「日本企業適合 × データ活用」で空白ポジションを獲得可能。ON24のCvent買収で日本市場投資が後退する今がチャンス。', 10, False, DARK_GREY),
    ])


# ==================================================================
# SLIDE 9: Positioning Map 1
# ==================================================================
def slide_09_positioning_map1():
    slide = add_slide()
    add_header(slide, 'Positioning Map 1: 機能深度 × 日本企業適合性',
               '右上象限（高機能×高適合）にネクプロを再配置する', 8)

    add_image(slide, 'positioning_map1.png', Inches(0.8), Inches(1.5), width=Inches(7.5))

    add_multiline_textbox(slide, Inches(8.5), Inches(2.0), Inches(4.2), Inches(4.5), [
        ('戦略意図', 14, True, NAVY),
        ('', 6, False, NAVY),
        ('現在：', 11, True, BLUE),
        ('機能は中程度、日本適合性は高い。', 10, False, DARK_GREY),
        ('しかしデータ分析・AI実行力でON24に劣後。', 10, False, DARK_GREY),
        ('', 6, False, NAVY),
        ('目標：', 11, True, ACCENT),
        ('エンゲージメントスコア＋AIコンテンツ生成で', 10, False, DARK_GREY),
        ('機能深度を引き上げ、右上象限を占有。', 10, False, DARK_GREY),
        ('', 6, False, NAVY),
        ('Zoom/ON24は日本適合性を短期に', 10, False, DARK_GREY),
        ('改善することが困難（構造的障壁）。', 10, False, DARK_GREY),
        ('', 8, False, NAVY),
        ('→ 「日本×高機能」は空白ポジション', 11, True, ACCENT),
    ])


# ==================================================================
# SLIDE 10: Positioning Map 2
# ==================================================================
def slide_10_positioning_map2():
    slide = add_slide()
    add_header(slide, 'Positioning Map 2: データ活用高度性 × 導入ハードル',
               'Sweet Spot = 高データ活用 × 低導入ハードル', 9)

    add_image(slide, 'positioning_map2.png', Inches(0.8), Inches(1.5), width=Inches(7.5))

    add_multiline_textbox(slide, Inches(8.5), Inches(2.0), Inches(4.2), Inches(4.5), [
        ('Sweet Spot戦略', 14, True, NAVY),
        ('', 6, False, NAVY),
        ('ON24は高データ活用だが', 10, False, DARK_GREY),
        ('導入ハードルも高い（英語・高額・海外対応）。', 10, False, DARK_GREY),
        ('', 6, False, NAVY),
        ('ネクプロの差別化：', 11, True, ACCENT),
        ('• テンプレートによる即時活用', 10, False, DARK_GREY),
        ('• 伴走型CSによるオンボーディング', 10, False, DARK_GREY),
        ('• Salesforceネイティブ連携', 10, False, DARK_GREY),
        ('• 日本語完全対応', 10, False, DARK_GREY),
        ('', 6, False, NAVY),
        ('→ データ活用を高めつつ', 11, True, ACCENT),
        ('  導入障壁を下げる「両立」が勝ち筋', 11, True, ACCENT),
    ])


# ==================================================================
# SLIDE 11: MECE Issues
# ==================================================================
def slide_11_mece():
    slide = add_slide()
    add_header(slide, '重要課題のMECE整理',
               '6視点×緊急度×インパクト×難易度で全社課題を構造化', 10)
    add_key_message_box(slide,
        '最優先は「プロダクトのデータ基盤化」と「解約率の改善」。この2つが全戦略の基盤。')

    data = [
        ['視点', '課題', '緊急度', 'インパクト', '難易度'],
        ['市場', 'AI時代の事業価値再定義の遅れ', '高', '大', '中'],
        ['顧客', '成果指標（商談化）までの接続が弱い', '高', '大', '中'],
        ['プロダクト', 'エンゲージメントスコアリング未実装', '高', '大', '中'],
        ['プロダクト', 'AIコンテンツ自動生成の未実装', '中', '中', '中'],
        ['GTM', '業種別提案テンプレートの不足', '高', '中', '低'],
        ['GTM', 'CSがコストセンターのまま', '中', '大', '中'],
        ['組織', 'PMM機能の不在', '高', '大', '中'],
        ['組織', '30名体制の人材ボトルネック', '高', '大', '高'],
        ['財務', '新収益柱の実行リスク', '高', '大', '高'],
        ['財務', 'LTV拡張余地（クロスセル不足）', '中', '大', '中'],
    ]
    make_table(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(4.5),
               11, 5, data, font_size=9,
               col_widths=[Inches(1.5), Inches(5.5), Inches(1.5), Inches(2.0), Inches(1.8)])


# ==================================================================
# SLIDE 12: Strategic Options
# ==================================================================
def slide_12_strategy_options():
    slide = add_slide()
    add_header(slide, '戦略オプション比較（3案）',
               '守るか、備えるか、攻めるか — 3つの道', 11)

    # 3 columns
    options = [
        ('梅：防衛型', '既存事業の効率最大化', GREY,
         ['投資: 現行水準維持', 'FY27: ¥740M（仮説）',
          'プロダクト: 既存改善のみ', 'GTM: 現行体制効率化',
          '組織: 変更なし',
          'リスク: 低（短期）→高（中長期）', 'リターン: 低（縮小均衡）']),
        ('竹：均衡型 ★推奨起点', '既存深耕＋データ機能拡張', BLUE,
         ['投資: +30-50%', 'FY27: ¥930M（仮説）',
          'プロダクト: スコアMVP+AI要約', 'GTM: 業種別PKG+CS高度化',
          '組織: PMM兼務設置',
          'リスク: 中（実行力分散）', 'リターン: 中（持続成長軌道）']),
        ('松：攻勢型', '基盤転換＋組織再編＋新収益', ACCENT,
         ['投資: +80-100%', 'FY27: ¥1,385M（計画値）',
          'プロダクト: フルスタック転換', 'GTM: 新セグメント+価格再設計',
          '組織: PMM/RevOps正式新設',
          'リスク: 高（キャッシュ・混乱）', 'リターン: 高（市場ポジション確立）']),
    ]

    for i, (title, subtitle, color, items) in enumerate(options):
        x = Inches(0.5 + i * 4.2)
        y = Inches(1.6)

        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      x, y, Inches(3.9), Inches(5.0))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0xF8, 0xF9, 0xFB)
        box.line.color.rgb = color
        box.line.width = Pt(2) if i == 1 else Pt(1)

        # Header
        hdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      x, y, Inches(3.9), Inches(0.7))
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = color
        hdr.line.fill.background()
        add_textbox(slide, x + Inches(0.15), y + Inches(0.02), Inches(3.6), Inches(0.35),
                    title, size=14, bold=True, color=WHITE)
        add_textbox(slide, x + Inches(0.15), y + Inches(0.35), Inches(3.6), Inches(0.3),
                    subtitle, size=9, color=RGBColor(0xDD, 0xDD, 0xDD))

        for j, item in enumerate(items):
            add_textbox(slide, x + Inches(0.2), y + Inches(0.85 + j * 0.48),
                        Inches(3.5), Inches(0.4),
                        f'• {item}', size=9, color=DARK_GREY)

    add_multiline_textbox(slide, Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.4), [
        ('推奨: ', 10, True, BLUE),
        ('「竹」案を6ヶ月で実行し、KPI達成を条件に「松」案へ段階移行', 10, True, NAVY),
    ])


# ==================================================================
# SLIDE 13: Recommended Strategy
# ==================================================================
def slide_13_recommended():
    slide = add_slide()
    add_header(slide, '推奨戦略：「竹」→「松」段階移行',
               '全力投球ではなく、実績で信頼を獲得しながらギアを上げる', 12)

    # 3 logic boxes
    logics = [
        ('Logic 1: キャッシュ耐性', NAVY,
         '30名・¥7.5億の企業が全方位投資は自殺行為。\n竹案なら+30-50%で既存キャッシュフロー維持可能。'),
        ('Logic 2: 小さく証明→大きく張る', BLUE,
         'エンゲージメントスコアMVPを6ヶ月で構築\n→ 5社で成果実証 → 本格投資判断。'),
        ('Logic 3: 組織能力の段階構築', LIGHT_BLUE,
         'いきなりPMM/RevOps正式新設は機能しない。\n兼務→成果確認→専任化のステップが現実的。'),
    ]
    for i, (title, color, desc) in enumerate(logics):
        x = Inches(0.5 + i * 4.2)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      x, Inches(1.6), Inches(3.9), Inches(1.8))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0xF5, 0xF7, 0xFA)
        box.line.color.rgb = color
        box.line.width = Pt(1.5)
        add_textbox(slide, x + Inches(0.2), Inches(1.7), Inches(3.5), Inches(0.3),
                    title, size=12, bold=True, color=color)
        add_textbox(slide, x + Inches(0.2), Inches(2.1), Inches(3.5), Inches(1.2),
                    desc, size=10, color=DARK_GREY)

    # Gate Review table
    add_textbox(slide, Inches(0.8), Inches(3.6), Inches(5), Inches(0.35),
                '6ヶ月後 Gate Review 判断基準', size=14, bold=True, color=NAVY)

    data = [
        ['KPI', 'Gate基準', '達成→', '未達→'],
        ['スコアMVP', '5社以上導入', '松案へ移行', '竹案延長・修正'],
        ['月次解約率', '1.3%以下', '松案へ移行', 'CS施策見直し'],
        ['ARPA（長期）', '¥160K以上', '松案へ移行', '価格再検討'],
        ['営業DX売上', '¥15M/半期以上', '投資拡大', 'ピボット検討'],
    ]
    make_table(slide, Inches(0.5), Inches(4.0), Inches(12.3), Inches(2.2),
               5, 4, data, font_size=10,
               col_widths=[Inches(2.5), Inches(3.3), Inches(3.3), Inches(3.2)])


# ==================================================================
# SLIDE 14: Product Initiatives
# ==================================================================
def slide_14_product():
    slide = add_slide()
    add_header(slide, '重点施策A：プロダクト（AI/データ基盤）',
               'エンゲージメントデータを「金鉱」から「製品」に変える', 13)

    data = [
        ['#', '施策', '時期', 'KPI', '責任'],
        ['A1', 'エンゲージメントスコアMVP\n視聴・チャット・Q&A・CTA統合', '0-6M', 'MVP完成→5社導入', 'プロダクト'],
        ['A2', 'Salesforce連携強化\nリードスコアリングへのネイティブフィード', '0-6M', 'SF連携ARPA +20%', 'プロダクト'],
        ['A3', '日本語AIコンテンツ生成\nウェビナー→要約・ブログ・クリップ', '6-12M', '利用率30%以上', 'プロダクト'],
        ['A4', 'API-firstアーキテクチャ移行', '6-18M', 'API化率80%', 'CTO'],
        ['A5', 'HubSpot/Marketo連携', '12-18M', '非SF新規20社/年', 'プロダクト'],
    ]
    make_table(slide, Inches(0.5), Inches(1.6), Inches(12.3), Inches(3.0),
               6, 5, data, font_size=9,
               col_widths=[Inches(0.5), Inches(5.5), Inches(1.3), Inches(3.0), Inches(2.0)])

    add_multiline_textbox(slide, Inches(0.8), Inches(4.8), Inches(11.5), Inches(2.0), [
        ('根拠（なぜこの順序か）', 12, True, NAVY),
        ('', 4, False, NAVY),
        ('• A1が最優先: ON24は参加者あたり40-50データポイントでインテントスコアを構築。AutodeskはパイプラインROI 12倍を達成。', 10, False, DARK_GREY),
        ('  日本市場ではまだ誰も本格提供していない → 先行者優位。', 10, False, DARK_GREY),
        ('', 4, False, NAVY),
        ('• A3は市場実証済み: Goldcastは年間13万件超のAI動画クリップを生成。日本語特化が差別化要因。', 10, False, DARK_GREY),
        ('', 4, False, NAVY),
        ('• A5は中期以降: 30名体制で3つのCRM/MA統合を同時開発は不可能。SF深化が先。', 10, False, DARK_GREY),
    ])


# ==================================================================
# SLIDE 15: GTM Initiatives
# ==================================================================
def slide_15_gtm():
    slide = add_slide()
    add_header(slide, '重点施策B：GTM（営業・CS・価格）',
               '「売り方」と「守り方」を同時に変える', 14)

    # Sales initiatives
    add_textbox(slide, Inches(0.8), Inches(1.5), Inches(5), Inches(0.35),
                '営業施策', size=14, bold=True, color=NAVY)
    s_data = [
        ['#', '施策', '時期', 'KPI'],
        ['B1', '業種別ソリューションPKG（4業種）', '0-6M', '成約率+5pt'],
        ['B2', '価格体系再設計（3層構造）', '6-12M', '新規ARPA ¥180K'],
        ['B3', '営業DX事業の立上げ加速', '0-12M', 'FY25 ¥33.4M'],
        ['B4', 'SF経由獲得チャネル構築', '0-6M', 'SF経由月2件'],
    ]
    make_table(slide, Inches(0.5), Inches(1.9), Inches(6.0), Inches(2.2),
               5, 4, s_data, font_size=9,
               col_widths=[Inches(0.4), Inches(3.0), Inches(0.9), Inches(1.7)])

    # CS initiatives
    add_textbox(slide, Inches(7.0), Inches(1.5), Inches(5), Inches(0.35),
                'CS施策', size=14, bold=True, color=BLUE)
    c_data = [
        ['#', '施策', '時期', 'KPI'],
        ['C1', 'オンボーディング標準化', '0-3M', '60日完了率90%'],
        ['C2', 'ヘルススコア導入', '0-6M', '予兆検知率70%'],
        ['C3', 'CS Profit Center化', '6-12M', 'CS起点¥30M/年'],
        ['C4', '戦略アカウント制（上位20社）', '0-6M', '上位NRR 120%'],
    ]
    make_table(slide, Inches(6.8), Inches(1.9), Inches(6.0), Inches(2.2),
               5, 4, c_data, font_size=9,
               col_widths=[Inches(0.4), Inches(3.0), Inches(0.9), Inches(1.7)])

    # Key insight
    add_multiline_textbox(slide, Inches(0.8), Inches(4.5), Inches(11.5), Inches(2.2), [
        ('CS Profit Center化の論理', 12, True, NAVY),
        ('', 4, False, NAVY),
        ('現状: CSは「解約防止」のコストセンター。成果が見えにくく投資判断が困難。', 10, False, DARK_GREY),
        ('', 4, False, NAVY),
        ('転換: 伴走支援を「成果連動型コンサルティング」に再定義。', 10, False, DARK_GREY),
        ('  • ウェビナー→商談化率の可視化（エンゲージメントスコア連動）', 10, False, DARK_GREY),
        ('  • 成果レポートの定期提供 → アップセル提案の自然な接点', 10, False, DARK_GREY),
        ('  • 目標: CS起点のアップセル売上 ¥30M/年', 10, False, DARK_GREY),
        ('', 4, False, NAVY),
        ('反証リスク: CS人員の「守り」工数を削ると解約率が悪化するリスク → 戦略アカウント制で優先度管理', 10, True, ACCENT_RED),
    ])


# ==================================================================
# SLIDE 16: Organization
# ==================================================================
def slide_16_organization():
    slide = add_slide()
    add_header(slide, '重点施策C：組織再編',
               '30名で勝つための「機能配置」と「採用計画」', 15)

    data = [
        ['#', '施策', '時期', 'KPI', '責任'],
        ['D1', 'PMM機能の兼務設置\n市場要件→製品要件の変換パイプライン', '0-3M', '月次市場要求レポート提出', 'CEO直轄'],
        ['D2', 'RevOps兼務設置\n営業/CS/プロダクトKPI統合', '3-6M', '統合ダッシュボード構築', 'COO相当'],
        ['D3', '戦略採用（3-5名）\nAIエンジニア1+PMM1+CS2+営業DX1', '6-12M', '採用充足率80%', '人事/CEO'],
        ['D4', 'KPIオーナー制度導入\n各KPIに個人名を紐付け', '0-1M', '全主要KPI\nオーナーアサイン完了', 'CEO'],
    ]
    make_table(slide, Inches(0.5), Inches(1.5), Inches(12.3), Inches(2.8),
               5, 5, data, font_size=9,
               col_widths=[Inches(0.5), Inches(5.0), Inches(1.3), Inches(3.3), Inches(2.2)])

    # Organization change visual
    add_textbox(slide, Inches(0.8), Inches(4.5), Inches(5), Inches(0.35),
                '組織変革のステップ', size=13, bold=True, color=NAVY)

    steps = [
        ('Month 1', 'KPIオーナー\nアサイン', NAVY),
        ('Month 1-3', 'PMM兼務\n設置', BLUE),
        ('Month 3-6', 'RevOps兼務\n設置', LIGHT_BLUE),
        ('Month 6-12', '戦略採用\n3-5名', ACCENT),
        ('Month 12+', '専任化\n判断', GREY),
    ]
    for i, (period, label, color) in enumerate(steps):
        x = Inches(0.8 + i * 2.4)
        y = Inches(5.0)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      x, y, Inches(2.0), Inches(1.3))
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()
        add_textbox(slide, x + Inches(0.1), y + Inches(0.1), Inches(1.8), Inches(0.25),
                    period, size=8, bold=True, color=WHITE)
        add_textbox(slide, x + Inches(0.1), y + Inches(0.4), Inches(1.8), Inches(0.8),
                    label, size=10, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

        if i < len(steps) - 1:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                            x + Inches(2.05), y + Inches(0.4),
                                            Inches(0.3), Inches(0.35))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = LIGHT_GREY
            arrow.line.fill.background()


# ==================================================================
# SLIDE 17: Roadmap
# ==================================================================
def slide_17_roadmap():
    slide = add_slide()
    add_header(slide, '3層実行ロードマップ（0-36ヶ月）',
               '止血 → 転換 → 成長', 16)

    add_image(slide, 'roadmap.png', Inches(0.3), Inches(1.5), width=Inches(12.5))


# ==================================================================
# SLIDE 18: KPI Tree
# ==================================================================
def slide_18_kpi_tree():
    slide = add_slide()
    add_header(slide, 'KPIツリーと経営モニタリング設計',
               '「測れないものは改善できない」— KPIオーナーを個人名で紐付け', 17)

    add_image(slide, 'kpi_tree.png', Inches(0.3), Inches(1.4), width=Inches(8))

    # Governance table
    add_textbox(slide, Inches(8.5), Inches(1.5), Inches(4), Inches(0.35),
                'ガバナンス体制', size=13, bold=True, color=NAVY)

    gov_data = [
        ['会議体', '頻度', 'アジェンダ'],
        ['経営KPI\nレビュー', '月次', 'KPI進捗・逸脱分析'],
        ['四半期戦略\nレビュー', '四半期', 'ロードマップ進捗\n方針修正'],
        ['Gate\nReview', 'M6,12,18', '戦略オプション\n移行判断'],
        ['スプリント\nレビュー', '隔週', '開発進捗\n優先度見直し'],
    ]
    make_table(slide, Inches(8.5), Inches(1.9), Inches(4.3), Inches(2.8),
               5, 3, gov_data, font_size=8,
               col_widths=[Inches(1.3), Inches(1.0), Inches(2.0)])

    # KPI Owner summary
    add_textbox(slide, Inches(8.5), Inches(5.0), Inches(4), Inches(0.35),
                'KPIオーナー（主要）', size=12, bold=True, color=NAVY)
    kpi_owners = [
        ['KPI', '目標(FY25)', 'オーナー'],
        ['ARR成長率', '+27%', 'CEO'],
        ['NRR', '100%+', 'CS責任者'],
        ['ARPA長期', '¥169K', '営業責任者'],
        ['解約率', '1.0%', 'CS責任者'],
        ['営業DX', '¥33.4M', 'DX責任者'],
    ]
    make_table(slide, Inches(8.5), Inches(5.4), Inches(4.3), Inches(1.5),
               6, 3, kpi_owners, font_size=8,
               col_widths=[Inches(1.5), Inches(1.3), Inches(1.5)])


# ==================================================================
# SLIDE 19: Decision Agenda
# ==================================================================
def slide_19_decision():
    slide = add_slide()
    add_header(slide, '本会議での決議依頼事項',
               '今日決めなければ、「何もしない」という選択を自動的にしたことになる', 18)

    decisions = [
        ('決議1', '投資配分の承認',
         'プロダクト45% / GTM35% / 組織20%の投資比率', NAVY),
        ('決議2', 'ポジショニング転換の正式承認',
         '「ウェビナーツール」→「B2Bエンゲージメント・インテリジェンス基盤」', BLUE),
        ('決議3', '組織再編の承認',
         'PMM兼務設置（即時）+ RevOps兼務（3ヶ月後）+ KPIオーナー制度', LIGHT_BLUE),
        ('決議4', '90日実行計画の承認',
         '短期施策パッケージの着手', ACCENT),
        ('決議5', '6ヶ月後Gate Reviewの設定',
         'KPI達成時の松案移行判断プロセス', GREY),
    ]

    for i, (label, title, desc, color) in enumerate(decisions):
        y = Inches(1.5 + i * 1.05)
        # Number circle
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                         Inches(0.8), y, Inches(0.55), Inches(0.55))
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()
        add_textbox(slide, Inches(0.82), y + Inches(0.05), Inches(0.55), Inches(0.45),
                    str(i + 1), size=18, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

        add_textbox(slide, Inches(1.6), y + Inches(0.0), Inches(5), Inches(0.35),
                    title, size=14, bold=True, color=NAVY)
        add_textbox(slide, Inches(1.6), y + Inches(0.35), Inches(8), Inches(0.3),
                    desc, size=10, color=DARK_GREY)

        # Status badge
        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        Inches(10.5), y + Inches(0.05), Inches(2.0), Inches(0.45))
        badge.fill.solid()
        badge.fill.fore_color.rgb = color
        badge.line.fill.background()
        add_textbox(slide, Inches(10.5), y + Inches(0.08), Inches(2.0), Inches(0.4),
                    '承認を推奨', size=10, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    # Next Steps
    add_textbox(slide, Inches(0.8), Inches(6.6), Inches(12), Inches(0.4),
                'Next Steps（承認後の即時アクション）', size=12, bold=True, color=NAVY)

    ns_items = 'KPIオーナーアサイン（1週間） → PMM兼務者選定（2週間） → スコアMVP要件定義KO（2週間） → 月次レビュー初回設定（1週間）'
    add_textbox(slide, Inches(0.8), Inches(6.95), Inches(12), Inches(0.35),
                ns_items, size=9, color=DARK_GREY)


# ==================================================================
# SLIDE 20: Appendix - Q&A
# ==================================================================
def slide_20_qa():
    slide = add_slide()
    add_header(slide, '想定Q&A（厳しめの指摘への回答）',
               'Appendix', 19)

    qa_data = [
        ['#', '想定質問', '回答要旨'],
        ['Q1', '30名でこの施策量を実行できるか？', '竹案で3施策に絞る。全部同時は不可。Gate Reviewで判断。'],
        ['Q2', '結局ウェビナーツールでは？', 'ON24実証：40-50データpt→ROI12倍。データ構造化が本質的価値。'],
        ['Q3', '営業DX ¥327M(FY27)は現実的か？', '¥15M/半期がリトマス試験。未達ならピボット。'],
        ['Q4', 'Zoom/Webex/Teamsに勝てるか？', '配信では勝てない。日本B2Bファネル最適化で差別化。'],
        ['Q5', 'SF依存は危険では？', '30名で3CRM同時開発は不可。SF深化→中期でAPI-first→拡張。'],
        ['Q6', '解約率1.0%は達成可能か？', 'FY22:3.6%→FY24:1.7%と改善中。ヘルススコア+標準化で可能。'],
        ['Q7', 'PMM/RevOps兼務で機能するか？', '3ヶ月で解像度UP→専任化判断。CEO直轄で工数20%確保。'],
        ['Q8', 'FY26-27計画は攻めすぎでは？', '竹案¥930Mがベースケース。営業DX未達時のダウンサイド準備要。'],
        ['Q9', 'bizibl/FanGrowthとの差別化は？', 'SF深度統合+データ構造化+Agentforce。ただし6-12ヶ月の猶予。'],
        ['Q10', 'ウェビナー市場自体が縮小しない？', '企業ウェビナー数:13→47回/年。ハイブリッド化で需要増大。'],
    ]
    make_table(slide, Inches(0.3), Inches(1.5), Inches(12.7), Inches(5.3),
               12, 3, qa_data, font_size=8,
               col_widths=[Inches(0.5), Inches(4.5), Inches(7.7)])


# ==================================================================
# Generate All Slides
# ==================================================================
def main():
    print('Generating slides...')
    slide_01_title()
    print('  1/20 Title')
    slide_02_exec_summary()
    print('  2/20 Executive Summary')
    slide_03_external()
    print('  3/20 External Environment')
    slide_04_current_revenue()
    print('  4/20 Current Revenue')
    slide_05_current_kpis()
    print('  5/20 Current KPIs')
    slide_06_accounts_new_rev()
    print('  6/20 Accounts & New Revenue')
    slide_07_swot()
    print('  7/20 SWOT')
    slide_08_competitive_table()
    print('  8/20 Competitive Table')
    slide_09_positioning_map1()
    print('  9/20 Positioning Map 1')
    slide_10_positioning_map2()
    print('  10/20 Positioning Map 2')
    slide_11_mece()
    print('  11/20 MECE Issues')
    slide_12_strategy_options()
    print('  12/20 Strategy Options')
    slide_13_recommended()
    print('  13/20 Recommended Strategy')
    slide_14_product()
    print('  14/20 Product Initiatives')
    slide_15_gtm()
    print('  15/20 GTM Initiatives')
    slide_16_organization()
    print('  16/20 Organization')
    slide_17_roadmap()
    print('  17/20 Roadmap')
    slide_18_kpi_tree()
    print('  18/20 KPI Tree')
    slide_19_decision()
    print('  19/20 Decision Agenda')
    slide_20_qa()
    print('  20/20 Q&A')

    prs.save(OUT_PATH)
    print(f'\nPresentation saved to: {OUT_PATH}')
    print(f'Total slides: {len(prs.slides)}')


if __name__ == '__main__':
    main()
